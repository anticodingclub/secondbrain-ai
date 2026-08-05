"""Upload through to extracted text, via the real ASGI app.

FastAPI's TestClient runs background tasks to completion before returning, so
by the time a request finishes here, parsing has already run — the same code
path production uses, just synchronously observable.
"""

from __future__ import annotations

import pytest
from httpx import AsyncClient

pytestmark = pytest.mark.integration

PREFIX = "/api/v1"

MARKDOWN = b"""# Deployment Guide

Some introduction.

## Docker Setup

The Dockerfile builds on python:3.11-slim and installs tesseract for OCR.
"""


async def sign_up(client: AsyncClient, email: str = "ada@example.com") -> str:
    response = await client.post(
        f"{PREFIX}/auth/register",
        json={"email": email, "password": "correct-horse-battery", "display_name": "Ada"},
    )
    return str(response.json()["access_token"])


def bearer(token: str) -> dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


async def upload(
    client: AsyncClient,
    token: str,
    *,
    filename: str = "guide.md",
    content: bytes = MARKDOWN,
):
    return await client.post(
        f"{PREFIX}/documents/upload",
        headers=bearer(token),
        files={"file": (filename, content, "text/markdown")},
    )


async def test_uploaded_document_is_parsed(client: AsyncClient) -> None:
    token = await sign_up(client)
    document_id = (await upload(client, token)).json()["document"]["id"]

    document = (await client.get(f"{PREFIX}/documents/{document_id}", headers=bearer(token))).json()

    assert document["status"] == "parsed"
    assert document["word_count"] > 0
    assert document["error_message"] is None


async def test_extracted_text_keeps_its_citation_anchors(client: AsyncClient) -> None:
    """The whole point of parsing at this layer: a citation must be able to
    say which section an answer came from."""
    token = await sign_up(client)
    document_id = (await upload(client, token)).json()["document"]["id"]

    response = await client.get(f"{PREFIX}/documents/{document_id}/text", headers=bearer(token))

    assert response.status_code == 200
    body = response.json()
    docker_block = next(b for b in body["blocks"] if "Dockerfile" in b["text"])
    assert docker_block["section_title"] == "Docker Setup"
    assert docker_block["heading_level"] == 2


async def test_page_numbers_survive_for_paged_formats(client: AsyncClient) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    token = await sign_up(client)

    presentation = Presentation()
    for index, heading in enumerate(["Architecture", "Roadmap"], start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = heading
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        box.text_frame.text = f"Body of slide {index}."

    import io

    buffer = io.BytesIO()
    presentation.save(buffer)

    response = await client.post(
        f"{PREFIX}/documents/upload",
        headers=bearer(token),
        files={
            "file": (
                "deck.pptx",
                buffer.getvalue(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    document_id = response.json()["document"]["id"]

    body = (
        await client.get(f"{PREFIX}/documents/{document_id}/text", headers=bearer(token))
    ).json()

    second = next(b for b in body["blocks"] if "slide 2" in b["text"])
    assert second["page_number"] == 2
    assert second["section_title"] == "Roadmap"
    assert body["page_count"] == 2


async def test_unreadable_file_is_marked_failed_not_deleted(client: AsyncClient) -> None:
    """A document that cannot be parsed is still the user's document. It must
    stay listed, with a reason, rather than vanishing."""
    token = await sign_up(client)
    document_id = (
        await client.post(
            f"{PREFIX}/documents/upload",
            headers=bearer(token),
            files={"file": ("broken.pdf", b"%PDF-1.7\nnot a real pdf body", "application/pdf")},
        )
    ).json()["document"]["id"]

    document = (await client.get(f"{PREFIX}/documents/{document_id}", headers=bearer(token))).json()

    assert document["status"] == "failed"
    assert document["error_message"]
    assert "broken.pdf" in document["error_message"]

    listing = await client.get(f"{PREFIX}/documents", headers=bearer(token))
    assert listing.json()["total"] == 1


async def test_scanned_pdf_reports_why_it_has_no_text(client: AsyncClient) -> None:
    """Without a warning, a scan is indistinguishable from a broken parser."""
    import io

    from pypdf import PdfWriter

    token = await sign_up(client)
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    buffer = io.BytesIO()
    writer.write(buffer)

    document_id = (
        await client.post(
            f"{PREFIX}/documents/upload",
            headers=bearer(token),
            files={"file": ("scan.pdf", buffer.getvalue(), "application/pdf")},
        )
    ).json()["document"]["id"]

    body = (
        await client.get(f"{PREFIX}/documents/{document_id}/text", headers=bearer(token))
    ).json()

    assert any("scanned" in warning.lower() for warning in body["warnings"])


async def test_duplicate_upload_does_not_reparse(client: AsyncClient) -> None:
    token = await sign_up(client)
    first = (await upload(client, token)).json()
    second = (await upload(client, token, filename="copy.md")).json()

    assert second["was_duplicate"] is True
    assert second["document"]["id"] == first["document"]["id"]
    # Still parsed from the first upload, not reset to pending by the second.
    assert second["document"]["status"] == "parsed"


async def test_text_endpoint_is_owner_scoped(client: AsyncClient) -> None:
    ada = await sign_up(client, "ada@example.com")
    grace = await sign_up(client, "grace@example.com")
    document_id = (await upload(client, ada)).json()["document"]["id"]

    response = await client.get(f"{PREFIX}/documents/{document_id}/text", headers=bearer(grace))

    assert response.status_code == 404


async def test_reparse_reruns_extraction(client: AsyncClient) -> None:
    token = await sign_up(client)
    document_id = (await upload(client, token)).json()["document"]["id"]

    response = await client.post(f"{PREFIX}/documents/{document_id}/reparse", headers=bearer(token))
    assert response.status_code == 200

    document = (await client.get(f"{PREFIX}/documents/{document_id}", headers=bearer(token))).json()
    assert document["status"] == "parsed"


async def test_deleting_a_document_removes_its_extracted_text(
    client: AsyncClient, settings
) -> None:
    """Extracted text is derived data and must not outlive its source."""
    token = await sign_up(client)
    document_id = (await upload(client, token)).json()["document"]["id"]

    before = list(settings.storage_path.rglob("*.blocks.json"))
    assert before, "parsing should have written an extracted-text artifact"

    await client.delete(f"{PREFIX}/documents/{document_id}", headers=bearer(token))

    assert list(settings.storage_path.rglob("*.blocks.json")) == []
    assert list(settings.storage_path.rglob("*")) == [
        path for path in settings.storage_path.rglob("*") if path.is_dir()
    ]


async def test_legacy_office_format_is_refused_with_advice(client: AsyncClient) -> None:
    token = await sign_up(client)

    response = await client.post(
        f"{PREFIX}/documents/upload",
        headers=bearer(token),
        files={"file": ("report.doc", b"\xd0\xcf\x11\xe0" + b"\x00" * 40, "application/msword")},
    )

    assert response.status_code == 415
    assert ".docx" in response.json()["message"]
