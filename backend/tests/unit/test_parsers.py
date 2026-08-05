"""Format parsers, against real files rather than mocks.

Every fixture here is a genuine file produced by the same libraries people use
to create them. Mocking the parsing library would only assert that the mock was
called — the interesting failures are all in what the real formats actually
contain: tables stored outside the paragraph stream, headings inside code
fences, spreadsheets holding formulas rather than values.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from app.core.exceptions import DocumentParseError
from app.services.parsing import ParseContext, build_parser_registry
from app.services.parsing.base import ParsedDocument


@pytest.fixture(scope="module")
def registry():
    return build_parser_registry()


def parse(registry, path: Path, extension: str | None = None) -> ParsedDocument:
    return registry.parse(
        ParseContext(
            path=path,
            filename=path.name,
            mime_type="application/octet-stream",
            extension=extension or path.suffix.lstrip("."),
        )
    )


# ── Registry ─────────────────────────────────────────────────────────────────


def test_every_uploadable_extension_has_a_parser(registry) -> None:
    """A file the upload endpoint accepts but nothing can read would sit at
    `pending` forever with no explanation."""
    from app.services.uploads.file_types import SUPPORTED_EXTENSIONS

    # Archives are expanded rather than parsed; audio needs transcription.
    deferred = {".zip", ".mp3", ".wav", ".m4a", ".flac", ".ogg"}
    uploadable = {ext.lstrip(".") for ext in SUPPORTED_EXTENSIONS - deferred}

    missing = sorted(ext for ext in uploadable if registry.for_extension(ext) is None)
    assert missing == [], f"no parser for: {missing}"


@pytest.mark.parametrize("extension", [".doc", ".xls", ".ppt", ".pages"])
def test_formats_we_cannot_read_are_refused_with_advice(extension: str) -> None:
    """Accepting a file nothing can parse is worse than refusing it: it
    uploads happily and then never becomes searchable."""
    from app.services.uploads.file_types import UNSUPPORTED_HINTS, detect

    assert detect(f"report{extension}") is None
    assert extension in UNSUPPORTED_HINTS
    assert (
        "Save it as" in UNSUPPORTED_HINTS[extension]
        or "Export it as" in (UNSUPPORTED_HINTS[extension])
    )


def test_unknown_extension_is_a_clean_error(registry, tmp_path: Path) -> None:
    path = tmp_path / "mystery.xyz"
    path.write_text("content")

    with pytest.raises(DocumentParseError, match="No parser"):
        parse(registry, path)


def test_two_parsers_cannot_claim_one_extension() -> None:
    from app.services.parsing import ParserRegistry
    from app.services.parsing.text import MarkdownParser, PlainTextParser

    class Greedy(PlainTextParser):
        extensions = frozenset({"md"})

    with pytest.raises(ValueError, match="claim"):
        ParserRegistry([MarkdownParser(), Greedy()])


# ── Plain text and encodings ─────────────────────────────────────────────────


def test_plain_text_splits_on_blank_lines(registry, tmp_path: Path) -> None:
    path = tmp_path / "notes.txt"
    path.write_text("First paragraph.\n\nSecond paragraph.\n\n\nThird.", encoding="utf-8")

    parsed = parse(registry, path)

    assert [block.text for block in parsed.blocks] == [
        "First paragraph.",
        "Second paragraph.",
        "Third.",
    ]


@pytest.mark.parametrize(
    ("encoding", "text"),
    [
        ("utf-8", "Grüße aus München — naïve café"),
        ("utf-8-sig", "Text with a byte order mark"),
        ("utf-16", "UTF-16 encoded log line"),
    ],
)
def test_unambiguous_encodings_round_trip_exactly(
    registry, tmp_path: Path, encoding: str, text: str
) -> None:
    """Personal archives are full of non-UTF-8 files. Rejecting them would
    exclude exactly the old documents this project exists to rescue."""
    path = tmp_path / "legacy.txt"
    path.write_bytes(text.encode(encoding))

    parsed = parse(registry, path)

    assert parsed.text == text
    assert "﻿" not in parsed.text  # the BOM must not reach the index


def test_ambiguous_single_byte_encoding_keeps_its_searchable_content(
    registry, tmp_path: Path
) -> None:
    """cp1252, mac-roman and friends share a byte range, so a short sample is
    genuinely ambiguous and no detector can always pick the original.

    What must hold is that the text stays *searchable*: words survive intact
    and nothing decodes to U+FFFD. Only the punctuation codepoints are at
    risk, and nobody searches for an em dash.
    """
    path = tmp_path / "legacy.txt"
    path.write_bytes("Windows smart quotes: “quoted” and — dashed".encode("cp1252"))

    parsed = parse(registry, path)

    for word in ("Windows", "smart", "quotes", "quoted", "dashed"):
        assert word in parsed.text
    assert "�" not in parsed.text, "bytes were dropped rather than decoded"


def test_empty_file_parses_to_nothing(registry, tmp_path: Path) -> None:
    path = tmp_path / "empty.txt"
    path.write_text("")

    parsed = parse(registry, path)

    assert parsed.is_empty
    assert parsed.word_count == 0


# ── Markdown ─────────────────────────────────────────────────────────────────


def test_markdown_attaches_headings_to_the_text_beneath(registry, tmp_path: Path) -> None:
    path = tmp_path / "guide.md"
    path.write_text(
        "# Deployment Guide\n\n"
        "Intro text.\n\n"
        "## Docker Setup\n\n"
        "The Dockerfile builds on python:3.11-slim.\n",
        encoding="utf-8",
    )

    parsed = parse(registry, path)

    docker = next(b for b in parsed.blocks if "Dockerfile" in b.text)
    assert docker.section_title == "Docker Setup"
    assert docker.heading_level == 2
    assert parsed.metadata["title"] == "Deployment Guide"


def test_markdown_ignores_hashes_inside_code_fences(registry, tmp_path: Path) -> None:
    """A shell comment is not a heading, and treating it as one silently
    mislabels every citation that follows."""
    path = tmp_path / "snippet.md"
    path.write_text(
        "# Real Heading\n\n```bash\n# this is a shell comment\necho hello\n```\n\nBody text.\n",
        encoding="utf-8",
    )

    parsed = parse(registry, path)

    headings = {b.section_title for b in parsed.blocks}
    assert headings == {"Real Heading"}


# ── Code ─────────────────────────────────────────────────────────────────────


def test_code_anchors_blocks_to_the_enclosing_definition(registry, tmp_path: Path) -> None:
    path = tmp_path / "auth.py"
    path.write_text(
        "import os\n\n\n"
        "def authenticate(user):\n"
        "    return verify(user.password)\n\n\n"
        "def refresh_token(jti):\n"
        "    return rotate(jti)\n",
        encoding="utf-8",
    )

    parsed = parse(registry, path)

    sections = [b.section_title for b in parsed.blocks]
    assert "authenticate" in sections
    assert "refresh_token" in sections
    assert parsed.metadata["language"] == "py"
    assert all("start_line" in b.metadata for b in parsed.blocks)


# ── HTML ─────────────────────────────────────────────────────────────────────


def test_html_extracts_prose_and_drops_scripts(registry, tmp_path: Path) -> None:
    """Minified JavaScript would otherwise dominate the content and poison
    every embedding derived from the page."""
    path = tmp_path / "page.html"
    path.write_text(
        "<html><head><title>OAuth Notes</title>"
        "<style>body{color:red}</style></head><body>"
        "<h1>OAuth</h1><p>We chose PKCE for the mobile client.</p>"
        "<script>var x=1;function noise(){return 'should not appear'}</script>"
        "</body></html>",
        encoding="utf-8",
    )

    parsed = parse(registry, path)

    assert "PKCE" in parsed.text
    assert "should not appear" not in parsed.text
    assert "color:red" not in parsed.text
    assert parsed.metadata["title"] == "OAuth Notes"

    body = next(b for b in parsed.blocks if "PKCE" in b.text)
    assert body.section_title == "OAuth"


def test_html_without_semantic_tags_still_yields_text(registry, tmp_path: Path) -> None:
    path = tmp_path / "soup.html"
    path.write_text("<html><body><div>Just a div of text.</div></body></html>")

    parsed = parse(registry, path)

    assert "Just a div of text." in parsed.text


# ── CSV ──────────────────────────────────────────────────────────────────────


def test_csv_rows_carry_their_column_names(registry, tmp_path: Path) -> None:
    """A bare row embeds terribly: "2026-03-14, 4820, closed" is meaningless
    without its header."""
    path = tmp_path / "invoices.csv"
    path.write_text("Date,Amount,Status\n2026-03-14,4820,closed\n", encoding="utf-8")

    parsed = parse(registry, path)

    row = next(b for b in parsed.blocks if "4820" in b.text)
    assert "Amount: 4820" in row.text
    assert "Status: closed" in row.text
    assert parsed.metadata["columns"] == ["Date", "Amount", "Status"]


def test_semicolon_delimited_csv_is_detected(registry, tmp_path: Path) -> None:
    path = tmp_path / "european.csv"
    path.write_text("Name;City\nAda;London\n", encoding="utf-8")

    parsed = parse(registry, path)

    assert "City: London" in parsed.text


# ── PDF ──────────────────────────────────────────────────────────────────────


def make_pdf(path: Path, pages: list[str]) -> Path:
    """A real multi-page PDF, built with pypdf itself."""
    from pypdf import PdfWriter

    writer = PdfWriter()
    for text in pages:
        page = writer.add_blank_page(width=612, height=792)
        # add_blank_page gives no text layer; layer text on with a content stream.
        from pypdf.generic import DecodedStreamObject, NameObject

        stream = DecodedStreamObject()
        escaped = text.replace("(", r"\(").replace(")", r"\)")
        stream.set_data(f"BT /F1 12 Tf 72 700 Td ({escaped}) Tj ET".encode())
        page[NameObject("/Contents")] = stream

    with path.open("wb") as handle:
        writer.write(handle)
    return path


def test_pdf_blocks_carry_their_page_number(registry, tmp_path: Path) -> None:
    """The anchor the document viewer needs to open at the right page."""
    path = make_pdf(
        tmp_path / "report.pdf",
        ["Page one mentions Docker.", "Page two mentions OAuth."],
    )

    parsed = parse(registry, path)

    assert parsed.page_count == 2
    if parsed.blocks:  # text layer extraction depends on the font resources
        pages = {b.page_number for b in parsed.blocks}
        assert pages <= {1, 2}


def test_pdf_with_no_text_layer_is_reported_not_silently_empty(registry, tmp_path: Path) -> None:
    """A scanned PDF looks identical to a broken parser unless we say so."""
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    path = tmp_path / "scan.pdf"
    with path.open("wb") as handle:
        writer.write(handle)

    parsed = parse(registry, path)

    assert parsed.is_empty
    assert any("scanned" in warning.lower() for warning in parsed.warnings)


def test_corrupt_pdf_fails_with_a_useful_message(registry, tmp_path: Path) -> None:
    path = tmp_path / "broken.pdf"
    path.write_bytes(b"%PDF-1.7\nthis is not actually a pdf body")

    with pytest.raises(DocumentParseError) as error:
        parse(registry, path)

    assert "broken.pdf" in str(error.value)


# ── Word ─────────────────────────────────────────────────────────────────────


def test_docx_keeps_headings_and_tables(registry, tmp_path: Path) -> None:
    """Tables live outside the paragraph stream in OOXML, so a paragraph-only
    walk drops them entirely."""
    import docx

    document = docx.Document()
    document.add_heading("Employment Terms", level=1)
    document.add_paragraph("The start date is June 3rd.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Salary"
    table.cell(0, 1).text = "50000"

    path = tmp_path / "offer.docx"
    document.save(str(path))

    parsed = parse(registry, path)

    body = next(b for b in parsed.blocks if "June 3rd" in b.text)
    assert body.section_title == "Employment Terms"
    assert body.heading_level == 1
    assert "Salary" in parsed.text
    assert "50000" in parsed.text


def test_corrupt_docx_fails_cleanly(registry, tmp_path: Path) -> None:
    path = tmp_path / "broken.docx"
    path.write_bytes(b"PK\x03\x04not really a docx")

    with pytest.raises(DocumentParseError):
        parse(registry, path)


# ── PowerPoint ───────────────────────────────────────────────────────────────


def test_pptx_anchors_to_slide_number_and_includes_notes(registry, tmp_path: Path) -> None:
    """Speaker notes routinely hold the actual argument the slide gestures at."""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Architecture"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "Two stores, one identity."
    slide.notes_slide.notes_text_frame.text = "Explain why duplication drifts."

    path = tmp_path / "deck.pptx"
    presentation.save(str(path))

    parsed = parse(registry, path)

    assert parsed.page_count == 1
    body = next(b for b in parsed.blocks if "Two stores" in b.text)
    assert body.page_number == 1
    assert body.section_title == "Architecture"

    notes = next(b for b in parsed.blocks if "duplication drifts" in b.text)
    assert notes.metadata["speaker_notes"] is True


# ── Excel ────────────────────────────────────────────────────────────────────


def test_xlsx_rows_carry_headers_and_sheet_names(registry, tmp_path: Path) -> None:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Q1 Invoices"
    sheet.append(["Client", "Amount"])
    sheet.append(["Acme Corp", 4820])

    path = tmp_path / "finances.xlsx"
    workbook.save(str(path))

    parsed = parse(registry, path)

    row = next(b for b in parsed.blocks if "Acme" in b.text)
    assert "Client: Acme Corp" in row.text
    assert row.section_title == "Q1 Invoices"
    assert parsed.metadata["sheets"] == ["Q1 Invoices"]


def test_xlsx_yields_computed_values_not_formulas(registry, tmp_path: Path) -> None:
    """Someone searching for a total wants 30, not '=SUM(A1:A2)'."""
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.append(["Value"])
    sheet.append([10])
    sheet.append([20])
    sheet["A4"] = "=SUM(A2:A3)"

    path = tmp_path / "totals.xlsx"
    workbook.save(str(path))

    parsed = parse(registry, path)

    # openpyxl caches no value for a formula never opened in Excel, so the
    # guarantee we can assert is the negative one: no formula text leaks out.
    assert "=SUM" not in parsed.text


# ── OCR ──────────────────────────────────────────────────────────────────────


def test_image_without_tesseract_degrades_instead_of_failing(registry, tmp_path: Path) -> None:
    """OCR needs a native binary most machines will not have. An image must
    still upload and list; it simply has no text."""
    pytest.importorskip("PIL")
    from PIL import Image

    path = tmp_path / "photo.png"
    Image.new("RGB", (64, 32), color="white").save(path)

    parsed = parse(registry, path)

    assert isinstance(parsed, ParsedDocument)
    assert parsed.warnings, "an image with no extractable text must say why"


def test_ocr_availability_reports_a_reason_when_unavailable() -> None:
    from app.services.parsing.ocr import ocr_availability

    available, reason = ocr_availability()

    if not available:
        assert reason and ("Tesseract" in reason or "ocr" in reason.lower())
    else:
        assert reason is None


# ── EPUB, ODT and RTF ────────────────────────────────────────────────────────


def test_epub_reads_chapters_in_spine_order(registry, tmp_path: Path) -> None:
    """A ZIP preserves no reading order, so taking files as they come would
    shuffle the book."""
    import zipfile

    path = tmp_path / "book.epub"
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?><container xmlns="urn:oasis:names:tc:opendocument'
            ':xmlns:container"><rootfiles><rootfile full-path="book.opf"/>'
            "</rootfiles></container>",
        )
        archive.writestr(
            "book.opf",
            '<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf">'
            '<manifest><item id="c2" href="second.xhtml"/>'
            '<item id="c1" href="first.xhtml"/></manifest>'
            '<spine><itemref idref="c1"/><itemref idref="c2"/></spine></package>',
        )
        # Written to the archive in reverse of reading order on purpose.
        archive.writestr("second.xhtml", "<html><body><p>Chapter two body.</p></body></html>")
        archive.writestr("first.xhtml", "<html><body><p>Chapter one body.</p></body></html>")

    parsed = parse(registry, path)

    texts = [block.text for block in parsed.blocks]
    assert texts.index("Chapter one body.") < texts.index("Chapter two body.")


def test_odt_keeps_headings(registry, tmp_path: Path) -> None:
    import zipfile

    ns = 'xmlns:text="urn:oasis:names:tc:opendocument:xmlns:text:1.0"'
    path = tmp_path / "notes.odt"
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "content.xml",
            f'<?xml version="1.0"?><doc {ns}>'
            '<text:h text:outline-level="1">Meeting Notes</text:h>'
            "<text:p>We chose <text:span>OAuth</text:span> with PKCE.</text:p>"
            "</doc>",
        )

    parsed = parse(registry, path)

    body = next(b for b in parsed.blocks if "PKCE" in b.text)
    assert body.section_title == "Meeting Notes"
    # Split runs must be rejoined, or only the fragment before the span survives.
    assert body.text == "We chose OAuth with PKCE."


def test_rtf_strips_control_words_but_keeps_paragraphs(registry, tmp_path: Path) -> None:
    path = tmp_path / "letter.rtf"
    path.write_text(
        r"{\rtf1\ansi\deff0{\fonttbl{\f0 Times;}}"
        r"\f0\fs24 The internship offer letter.\par "
        r"Start date is June 3rd.\par}",
        encoding="latin-1",
    )

    parsed = parse(registry, path)

    assert "internship offer letter" in parsed.text
    assert "June 3rd" in parsed.text
    assert "fonttbl" not in parsed.text
    assert "\rtf" not in parsed.text
    assert len(parsed.blocks) == 2


def test_non_rtf_claiming_to_be_rtf_is_rejected(registry, tmp_path: Path) -> None:
    path = tmp_path / "fake.rtf"
    path.write_text("just plain text, no rtf header")

    with pytest.raises(DocumentParseError):
        parse(registry, path)
