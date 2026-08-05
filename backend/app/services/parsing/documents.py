"""PDF, Word, PowerPoint and Excel."""

from __future__ import annotations

from typing import Any

from app.core.exceptions import DocumentParseError
from app.core.logging import get_logger
from app.services.parsing.base import (
    DocumentParser,
    ParseContext,
    ParsedDocument,
    TextBlock,
)
from app.services.parsing.text import MAX_BLOCK_CHARS, split_paragraphs

logger = get_logger(__name__)


class PdfParser(DocumentParser):
    """PDF text, one block per paragraph, anchored to its page.

    Page numbers are the whole reason this is not just `pdftotext`: the
    document viewer in Phase 15 opens directly to a cited page, and that is
    only possible if the page travels with the text from here.

    Pages that yield no text are reported rather than silently dropped, since
    a PDF of scans looks identical to a broken parse unless someone says so.
    """

    extensions = frozenset({"pdf"})

    def parse(self, context: ParseContext) -> ParsedDocument:
        from pypdf import PdfReader
        from pypdf.errors import PyPdfError

        try:
            reader = PdfReader(str(context.path))
        except PyPdfError as exc:
            raise DocumentParseError(f"{context.filename!r} is not a readable PDF.") from exc

        if reader.is_encrypted:
            # An empty user password is common for "print-protected" files and
            # costs nothing to try.
            try:
                if not reader.decrypt(""):
                    raise DocumentParseError(f"{context.filename!r} is password-protected.")
            except (PyPdfError, NotImplementedError) as exc:
                raise DocumentParseError(f"{context.filename!r} is password-protected.") from exc

        blocks: list[TextBlock] = []
        empty_pages: list[int] = []

        for index, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception as exc:
                # One malformed page must not lose the other 199.
                logger.warning("pdf_page_failed", page=index, error=str(exc))
                empty_pages.append(index)
                continue

            if not text.strip():
                empty_pages.append(index)
                continue

            blocks.extend(
                TextBlock(text=paragraph, page_number=index) for paragraph in split_paragraphs(text)
            )

        page_count = len(reader.pages)
        warnings: list[str] = []
        if empty_pages and len(empty_pages) == page_count:
            warnings.append("No text layer found. This looks like a scanned PDF — OCR is needed.")
        elif empty_pages:
            warnings.append(f"{len(empty_pages)} of {page_count} pages had no text.")

        return ParsedDocument(
            blocks=blocks,
            page_count=page_count,
            metadata=_pdf_metadata(reader),
            warnings=tuple(warnings),
        )


def _pdf_metadata(reader: Any) -> dict[str, Any]:
    try:
        info = reader.metadata or {}
    except Exception:
        return {}

    fields = {
        "title": info.get("/Title"),
        "author": info.get("/Author"),
        "subject": info.get("/Subject"),
        "creator": info.get("/Creator"),
    }
    return {key: str(value) for key, value in fields.items() if value}


class DocxParser(DocumentParser):
    """Word documents, keeping heading structure and tables."""

    extensions = frozenset({"docx"})

    def parse(self, context: ParseContext) -> ParsedDocument:
        import docx
        from docx.opc.exceptions import PackageNotFoundError

        try:
            document = docx.Document(str(context.path))
        except (PackageNotFoundError, KeyError, ValueError) as exc:
            raise DocumentParseError(
                f"{context.filename!r} is not a readable Word document."
            ) from exc

        blocks: list[TextBlock] = []
        section: str | None = None
        level: int | None = None

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            style = (paragraph.style.name or "") if paragraph.style else ""
            if style.startswith("Heading"):
                level = _heading_level(style)
                section = text
                blocks.append(TextBlock(text=text, section_title=section, heading_level=level))
            elif style == "Title":
                level, section = 1, text
                blocks.append(TextBlock(text=text, section_title=section, heading_level=1))
            else:
                blocks.append(
                    TextBlock(
                        text=text[:MAX_BLOCK_CHARS],
                        section_title=section,
                        heading_level=level,
                    )
                )

        # Tables are separate from paragraphs in the OOXML model, so a
        # paragraph-only walk silently drops every table in the file.
        for number, table in enumerate(document.tables, start=1):
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    blocks.append(
                        TextBlock(
                            text=" | ".join(cells),
                            section_title=section,
                            metadata={"table": number},
                        )
                    )

        properties = document.core_properties
        metadata = {
            key: value
            for key, value in {
                "title": properties.title,
                "author": properties.author,
                "subject": properties.subject,
            }.items()
            if value
        }
        return ParsedDocument(blocks=blocks, metadata=metadata)


def _heading_level(style_name: str) -> int:
    tail = style_name.removeprefix("Heading").strip()
    return int(tail) if tail.isdigit() else 1


class PptxParser(DocumentParser):
    """Slides, anchored by slide number, including speaker notes."""

    extensions = frozenset({"pptx"})

    def parse(self, context: ParseContext) -> ParsedDocument:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError

        try:
            presentation = Presentation(str(context.path))
        except (PackageNotFoundError, KeyError, ValueError) as exc:
            raise DocumentParseError(
                f"{context.filename!r} is not a readable presentation."
            ) from exc

        blocks: list[TextBlock] = []

        for number, slide in enumerate(presentation.slides, start=1):
            title = _slide_title(slide)

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        blocks.append(
                            TextBlock(
                                text=text[:MAX_BLOCK_CHARS],
                                page_number=number,
                                section_title=title,
                            )
                        )

            # Notes routinely hold the actual argument the slide only gestures
            # at, so they are worth indexing and worth labelling as notes.
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    blocks.append(
                        TextBlock(
                            text=notes[:MAX_BLOCK_CHARS],
                            page_number=number,
                            section_title=title,
                            metadata={"speaker_notes": True},
                        )
                    )

        return ParsedDocument(
            blocks=blocks,
            page_count=len(presentation.slides),
            metadata={"slide_count": len(presentation.slides)},
        )


def _slide_title(slide: Any) -> str | None:
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            return str(slide.shapes.title.text).strip()
    except (AttributeError, ValueError):
        pass
    return None


class XlsxParser(DocumentParser):
    """Spreadsheets, one block per row, anchored by sheet.

    `read_only` and `data_only` matter here: the first streams rather than
    building the whole workbook in memory, and the second yields cached
    computed values instead of `=SUM(B2:B40)`, which is what a person
    searching for a number actually wants.
    """

    extensions = frozenset({"xlsx", "xlsm"})

    MAX_ROWS_PER_SHEET = 5000

    def parse(self, context: ParseContext) -> ParsedDocument:
        import openpyxl
        from openpyxl.utils.exceptions import InvalidFileException

        try:
            workbook = openpyxl.load_workbook(str(context.path), read_only=True, data_only=True)
        except (InvalidFileException, KeyError, ValueError) as exc:
            raise DocumentParseError(
                f"{context.filename!r} is not a readable spreadsheet."
            ) from exc

        blocks: list[TextBlock] = []
        warnings: list[str] = []

        try:
            for index, sheet in enumerate(workbook.worksheets, start=1):
                header: list[str] = []
                rows_seen = 0

                for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    values = ["" if value is None else str(value).strip() for value in row]
                    if not any(values):
                        continue

                    if not header:
                        header = values
                        blocks.append(
                            TextBlock(
                                text=f"{sheet.title}: " + " | ".join(v for v in values if v),
                                page_number=index,
                                section_title=sheet.title,
                            )
                        )
                        continue

                    rows_seen += 1
                    if rows_seen > self.MAX_ROWS_PER_SHEET:
                        warnings.append(
                            f"Sheet {sheet.title!r} was truncated at "
                            f"{self.MAX_ROWS_PER_SHEET} rows."
                        )
                        break

                    pairs = [
                        f"{name}: {value}"
                        for name, value in zip(header, values, strict=False)
                        if value and name
                    ]
                    if pairs:
                        blocks.append(
                            TextBlock(
                                text=" | ".join(pairs)[:MAX_BLOCK_CHARS],
                                page_number=index,
                                section_title=sheet.title,
                                metadata={"row_number": row_number},
                            )
                        )
        finally:
            # read_only mode holds an open file handle until closed.
            workbook.close()

        sheet_names = [sheet.title for sheet in workbook.worksheets]
        return ParsedDocument(
            blocks=blocks,
            page_count=len(sheet_names),
            metadata={"sheets": sheet_names},
            warnings=tuple(warnings),
        )
